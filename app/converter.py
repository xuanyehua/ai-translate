import io
import re
from pathlib import Path

import mistune


def _parse_md_blocks(markdown: str) -> list[dict]:
    """将 Markdown 按段落拆分为带类型的 block 列表。"""
    blocks = markdown.split("\n\n")
    result = []
    for b in blocks:
        b = b.strip()
        if not b:
            continue
        if b.startswith("#"):
            result.append({"type": "heading", "text": b})
        elif b.startswith("- ") or b.startswith("* "):
            result.append({"type": "list", "text": b})
        elif b.startswith("```"):
            result.append({"type": "code", "text": b})
        elif b.startswith(">"):
            result.append({"type": "quote", "text": b})
        elif re.match(r"!\[.*\]\(.*\)", b):
            result.append({"type": "image", "text": b})
        else:
            result.append({"type": "paragraph", "text": b})
    return result


def convert_to_docx(markdown: str) -> bytes:
    """将翻译后的 Markdown 转为 .docx 文件字节流。"""
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    blocks = _parse_md_blocks(markdown)
    for block in blocks:
        text = re.sub(r"[#*>`\-\[\]!()]", "", block["text"]).strip()
        if not text:
            continue
        if block["type"] == "heading":
            level = len(block["text"]) - len(block["text"].lstrip("#"))
            doc.add_heading(text, level=min(level, 3))
        elif block["type"] in ("list",):
            p = doc.add_paragraph(text, style="List Bullet")
        elif block["type"] == "code":
            p = doc.add_paragraph()
            run = p.add_run(text)
            run.font.name = "Courier New"
            run.font.size = Pt(9)
        else:
            doc.add_paragraph(text)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def convert_to_pdf(markdown: str) -> bytes:
    """将翻译后的 Markdown 转为 .pdf 文件字节流（基础版本）。"""
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)

    blocks = _parse_md_blocks(markdown)
    for block in blocks:
        text = re.sub(r"[#*>`\-\[\]!()]", "", block["text"]).strip()
        if not text:
            continue
        if block["type"] == "heading":
            level = len(block["text"]) - len(block["text"].lstrip("#"))
            pdf.set_font("Helvetica", style="B", size=18 - level * 2)
            pdf.multi_cell(0, 10, text)
            pdf.ln(2)
        else:
            pdf.set_font("Helvetica", size=12)
            pdf.multi_cell(0, 7, text)
            pdf.ln(2)

    return pdf.output()


def convert_to_pptx(markdown: str) -> bytes:
    """将翻译后的 Markdown 转为 .pptx 文件字节流（每页标题 + 内容）。"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Split by ## headings to create slides
    sections = markdown.split("\n## ")
    for i, section in enumerate(sections):
        lines = section.strip().split("\n")
        if not lines:
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

        # Title
        title = lines[0].lstrip("#").strip()
        left = Inches(1)
        top = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, Inches(11.333), Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Inches(0.6)
        p.font.bold = True

        # Content
        body_lines = lines[1:] if len(lines) > 1 else []
        clean_lines = [re.sub(r"^[#*>\-\s]+", "", l).strip() for l in body_lines]
        clean_lines = [l for l in clean_lines if l]
        if clean_lines:
            top = Inches(1.8)
            txBox = slide.shapes.add_textbox(left, top, Inches(11.333), Inches(5))
            tf = txBox.text_frame
            tf.word_wrap = True
            for j, line in enumerate(clean_lines):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Inches(0.35)
                p.space_after = Inches(0.15)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# Format -> (output extension, MIME type)
_OUTPUT_MAP: dict[str, tuple[str, str]] = {
    "md": ("md", "text/markdown"),
    "docx": ("docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
    "pptx": ("pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    "pdf": ("pdf", "application/pdf"),
    # image formats + xlsx → output as markdown
}

_CONVERTERS = {
    "docx": convert_to_docx,
    "pptx": convert_to_pptx,
    "pdf": convert_to_pdf,
}


def convert(markdown: str, ext: str) -> tuple[bytes, str, str]:
    """
    将翻译后的 Markdown 转为原文档格式。
    返回 (文件字节流, MIME type, 输出扩展名)。
    """
    converter = _CONVERTERS.get(ext)
    out_ext, mime = _OUTPUT_MAP.get(ext, ("md", "text/markdown"))

    if converter:
        return converter(markdown), mime, out_ext
    else:
        return markdown.encode("utf-8"), mime, out_ext
